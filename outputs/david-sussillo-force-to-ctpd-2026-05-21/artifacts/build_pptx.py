# -*- coding: utf-8 -*-
"""Build the balanced 35-slide editable PPTX (same axis as the HTML deck).
Reproducible: reads images/, emits output/deck.pptx.
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = pathlib.Path(__file__).resolve().parent
IMG = HERE.parent / "images"

BG   = RGBColor(0x11,0x10,0x0e)
BG2  = RGBColor(0x1a,0x18,0x15)
FG   = RGBColor(0xfa,0xf9,0xf5)
MUT  = RGBColor(0xc2,0xc0,0xb6)
GREEN= RGBColor(0x5d,0xca,0xa5)
PURP = RGBColor(0xaf,0xa9,0xec)
BROWN= RGBColor(0xf0,0x99,0x7b)
GOLD = RGBColor(0xef,0x9f,0x27)
ACC = {"green":GREEN,"purple":PURP,"brown":BROWN,"gold":GOLD,None:GOLD}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def new(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0,0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def tb(s, x, y, w, h):
    box = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    return tf

def setp(p, text, size, color=FG, bold=False, align=PP_ALIGN.LEFT):
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    f.name = "Microsoft JhengHei"
    return p

def title(s, text, kicker=None, kc=None):
    if kicker:
        tf = tb(s, 0.7, 0.45, 12, 0.6)
        setp(tf.paragraphs[0], kicker, 14, ACC[kc], bold=True)
    tf = tb(s, 0.7, 0.95, 12, 1.2)
    setp(tf.paragraphs[0], text, 34, FG, bold=True)

def bullets(s, items, x=0.9, y=2.4, w=11.5, size=22, gap=True):
    tf = tb(s, x, y, w, 4.5)
    for i,it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        bold = it.startswith("**")
        txt = it.replace("**","")
        setp(p, ("•  " if not it.startswith("  ") else "    – ")+txt.strip(), size, FG if not bold else FG, bold=bold)
        p.space_after = Pt(10 if gap else 4)

def center_big(s, lines, base=40, sub=None):
    tf = tb(s, 1, 2.4, 11.3, 3.2); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i,l in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        setp(p, l, base, FG, bold=True, align=PP_ALIGN.CENTER); p.space_after=Pt(8)
    if sub:
        p = tf.add_paragraph(); setp(p, sub, 18, MUT, align=PP_ALIGN.CENTER)

def idea(s, lines, color="purple", y=2.6, size=26):
    bar = s.shapes.add_shape(1, Inches(1.1),Inches(y), Inches(0.1),Inches(1.7))
    bar.fill.solid(); bar.fill.fore_color.rgb=ACC[color]; bar.line.fill.background(); bar.shadow.inherit=False
    tf = tb(s, 1.4, y, 10.6, 1.9); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,l in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        setp(p, l, size, FG, bold=True); p.space_after=Pt(6)

def note(s, text, y=6.2):
    tf = tb(s, 0.9, y, 11.5, 1.0)
    setp(tf.paragraphs[0], text, 15, MUT)

def math(s, lines, y=2.6, size=28):
    tf = tb(s, 1, y, 11.3, 3); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,l in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        setp(p, l, size, GREEN if i%2==0 else FG, align=PP_ALIGN.CENTER); p.space_after=Pt(14)

def cards(s, items, y=2.6, h=2.6):
    n=len(items); gap=0.4; tot=12.0; cw=(tot-gap*(n-1))/n
    for i,(ttl,body,col) in enumerate(items):
        x=0.7+i*(cw+gap)
        box=s.shapes.add_shape(1, Inches(x),Inches(y),Inches(cw),Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb=BG2; box.line.color.rgb=ACC[col]; box.line.width=Pt(1.2); box.shadow.inherit=False
        tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.2); tf.margin_right=Inches(0.2)
        setp(tf.paragraphs[0], ttl, 19, ACC[col], bold=True)
        p=tf.add_paragraph(); setp(p, body, 16, FG); p.space_before=Pt(8)

def image(s, name, cap=None, top=1.9, maxh=4.4):
    p = IMG/name
    pic = s.shapes.add_picture(str(p), 0, Inches(top))
    # scale to maxh
    ratio = pic.width/pic.height
    nh = Inches(maxh); nw = int(nh*ratio)
    if nw > Inches(11.5):
        nw = Inches(11.5); nh = int(nw/ratio)
    pic.height=nh; pic.width=nw
    pic.left = int((W-nw)/2); pic.top=Inches(top)
    if cap:
        tf=tb(s, 0.7, top+maxh+0.15, 12, 0.7)
        setp(tf.paragraphs[0], cap, 14, MUT, align=PP_ALIGN.CENTER)

def vote(s, q, opts, ans, reveal):
    title(s, q, kicker="Check-Point", kc="gold")
    tf=tb(s, 1.2, 2.5, 10.5, 3)
    for i,o in enumerate(opts):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        is_ans = o.startswith(ans+".")
        setp(p, o, 20, GREEN if is_ans else FG, bold=is_ans); p.space_after=Pt(12)
    note(s, "（討論題）"+reveal, y=6.0)

def pivot_card(s, x, y, w, h, label, frm, to, col, src):
    box=s.shapes.add_shape(1, Inches(x),Inches(y),Inches(w),Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb=BG2; box.line.color.rgb=ACC[col]; box.line.width=Pt(1.5); box.shadow.inherit=False
    tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.25)
    setp(tf.paragraphs[0], label, 17, ACC[col], bold=True)
    p=tf.add_paragraph(); setp(p, frm+"  →  "+to, 16, FG); p.space_before=Pt(6)
    p=tf.add_paragraph(); setp(p, src, 13, MUT); p.space_before=Pt(4)

# ===================== 35 SLIDES =====================
# 1 cover
s=new();
tf=tb(s,1,0.5,11.3,0.6); setp(tf.paragraphs[0],"COMPUTATIONAL NEUROSCIENCE · 60 MIN",14,GREEN,bold=True,align=PP_ALIGN.CENTER)
tf=tb(s,1,1.9,11.3,2.4); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
setp(tf.paragraphs[0],"從 FORCE 到 Population Dynamics",44,FG,bold=True,align=PP_ALIGN.CENTER)
p=tf.add_paragraph(); setp(p,"一個演算法如何翻轉我們理解大腦的方式",24,MUT,align=PP_ALIGN.CENTER)
idea(s,["FORCE 不是演算法升級，而是讓 RNN 從「無法馴服的混沌」",
        "變成「可解剖的動力系統」——這個翻轉，生出了 CTPD 框架。"],"purple",y=4.4,size=20)
note(s,"重讀 David Sussillo 研究主線 · 並親手把它跑一遍", y=6.6)

# 2 hook
s=new(); title(s,"Hook：1000 個神經元的沉默")
cards(s,[("我們能記錄的","同時記錄上千個神經元：array、雙光子、Neuropixels…","green"),
         ("我們能理解的","把 tuning curve 一條條畫出，面對前額葉高維時間動態——卡住。","brown")])
idea(s,["記錄變多，理解沒等比例變多。","single-neuron tuning curve 不夠用了。"],"brown",y=5.3,size=22)

# 3 one big idea
s=new(); title(s,"本場的 One Big Idea")
pivot_card(s,0.7,2.6,3.7,1.6,"工具","混沌","可訓練","green","FORCE 2009")
pivot_card(s,4.8,2.6,3.7,1.6,"動力系統","黑盒","可解剖","purple","Black Box 2013")
pivot_card(s,8.9,2.6,3.7,1.6,"框架","方法","典範","gold","CTPD 2020")
center_big(s,["今天就講清楚這條翻轉是怎麼發生的。"],base=24)

# 4 roadmap of talk
s=new(); title(s,"今天的路線")
bullets(s,["**1  FORCE 的數學核心 — 為何誤差從頭就小",
           "**2  樞紐轉折 — 打開黑盒",
           "**3  四個轉折路線圖",
           "**4  親手重現 — FORCE 模擬 + fixed-point 分析",
           "**5  收斂到 CTPD + 批判"], size=24)
note(s,"★ 含三個親手跑的 numpy 實驗（零捏造）")

# 5 dilemma
s=new(); title(s,"混沌 RNN 的困境", kicker="PART 1", kc="green")
center_big(s,["把輸出誤差回饋進遞迴網路 → 回饋反覆放大（reverberating）→ 訓練發散"],base=22)
cards(s,[("前人（echo-state）","刻意讓網路無輸入時靜止才能訓 → 犧牲混沌動態。",None),
         ("FORCE","保留混沌，在訓練中主動壓制它。","green")], y=4.2, h=2.2)

# 6 pivot1
s=new(); title(s,"轉折 1：混沌 bug → feature")
idea(s,["FORCE 保留混沌的自發活動，","用強而快的突觸修改在訓練早期壓制它。"],"green")
note(s,"關鍵：回饋訊號 ≈ f(t)，振幅夠時可誘導 chaotic → nonchaotic transition。", y=5.0)

# 7 eq1-3
s=new(); title(s,"核心設定（式 1–3）")
math(s,["z(t) = wᵀ r(t)              … (1)",
        "e₋(t) = wᵀ(t−Δt) r(t) − f(t)   … (2)",
        "e₊(t) = wᵀ(t) r(t) − f(t)      … (3)"], size=24)
note(s,"輸出是 firing rate 的線性讀出；分別記更新「前」與「後」的誤差。")

# 8 CP1
s=new(); vote(s,"FORCE 與一般 gradient descent 最大差異是？",
  ["A. 用更大的 batch","B. 學習率調更小","C. 誤差全程維持很小，而非由大慢慢縮小","D. 不需要 target"],
  "C","答案 C：FORCE 不走「減小大誤差」，而是「維持小誤差」。")

# 9 eq4-6
s=new(); title(s,"RLS 更新式與 P 矩陣（式 4–6）")
math(s,["w(t) = w(t−Δt) − e₋(t) P(t) r(t)        … (4)",
        "P(t) = P − (P r rᵀ P)/(1 + rᵀ P r),  P(0)=I/α   … (5,6)"], size=22)
cards(s,[("P ≈ 相關矩陣的逆","= 二階（Newton-like）更新","purple"),
         ("每個突觸有自己","隨資料調整的有效步長（非純量學習率）",None)], y=4.6, h=1.9)

# 10 eq7-8 why small
s=new(); title(s,"為何誤差從頭就小（式 7–8）")
math(s,["e₋(Δt) = −α f / (α + rᵀr),   rᵀr ~ O(N)",
        "e₊(t) = e₋(t) · (1 − rᵀ P r)"], y=2.3, size=24)
bullets(s,["**α≪N → 第一步誤差 ≈ −αf/N，極小",
           "**rᵀPr：初期≈1（e₊≈0）→ 後期→0（收斂）",
           "誤差全程小 → 回饋擾動小 → 混沌被壓住"], y=4.4, size=20)

# 11 black box
s=new(); title(s,"訓練好了，然後呢？", kicker="PART 2", kc="purple")
center_big(s,["輸入  →   ?   →  正確輸出"],base=36)
idea(s,["監督式訓練規定了「做什麼」，卻沒規定「怎麼做」。","網路成了黑盒。"],"purple",y=4.6,size=22)

# 12 pivot2 key
s=new(); title(s,"關鍵洞見（轉折 2）★")
idea(s,["「怎麼做」就是一個可研究的科學問題。","把訓練好的 RNN 當非線性動力系統 (NLDS) 來逆向工程。"],"purple",size=24)
note(s,"Sussillo & Barak 2013 · Opening the Black Box · 整條演化的支點。", y=5.2)

# 13 ds language
s=new(); title(s,"動力系統語言入門")
cards(s,[("State space","每軸一個神經元",None),("Trajectory","狀態隨時間的路徑",None),
         ("Flow field","每點往哪走",None),("Fixed/slow point","速度≈0 的點","purple")])
note(s,"鐘擺（位置×速度）→ 神經群體（N 維狀態）的同一套語言。")

# 14 fixed/slow
s=new(); title(s,"Fixed points 與 slow points")
bullets(s,["在相空間找速度近乎為零的點",
           "對該點做線性化（Jacobian）",
           "特徵值 → 穩定 / 不穩定方向",
           "由幾何結構反推網路如何運算"], size=22)
note(s,"逆向工程三例：3-bit flip-flop、sine generator、2-point average — 機制皆可被讀出。")

# 15 pivot point why
s=new(); title(s,"這一步為何是樞紐")
idea(s,["從「會動」到「能讀懂」。"],"purple",size=30,y=2.6)
center_big(s,["沒有它，FORCE 只是工程技巧；","有了它，RNN 成為能生成假設的動力系統模型。"],base=22,)
note(s,"Sussillo & Barak 2013, Neural Computation — 後續一切的支點。")

# 16 roadmap full
s=new(); title(s,"四個轉折：FORCE → CTPD 全景", kicker="PART 3", kc=None)
pivot_card(s,0.7,1.9,12,1.05,"轉折 1 ｜ 混沌 bug→feature","保留混沌","主動壓制","green","FORCE 2009")
pivot_card(s,0.7,3.05,12,1.05,"轉折 2 ｜ 本體論翻轉 ★","黑盒工具","可逆向工程的 NLDS","purple","Black Box 2013")
pivot_card(s,0.7,4.2,12,1.05,"轉折 3 ｜ 模型→大腦","RNN 語言","真實神經群體","brown","Barak/Romo 2013 · Sussillo 2014")
pivot_card(s,0.7,5.35,12,1.05,"轉折 4 ｜ input 成為重構者","訊號","x₀+u(t) 指定運算","purple","Driscoll/Golub/Sussillo 2018")

# 17 CP2
s=new(); vote(s,"你認為哪個轉折最關鍵？",
  ["A. 混沌 bug→feature","B. 黑盒→可解剖的動力系統","C. 模型→大腦","D. input 成為重構者"],
  "B","講者觀點 B：它是支點，後三者才得以展開。")

# 18 pivot3
s=new(); title(s,"轉折 3：分析模型 → 分析大腦")
idea(s,["RNN 的動力系統語言（state space、trajectory、attractor）","反向套回真實神經群體資料。"],"brown")
note(s,"顛覆 single-neuron tuning 典範：要懂單一神經元，先懂 population。 · Barak/Romo 2013 · Sussillo 2014", y=5.2)

# 19 pivot4
s=new(); title(s,"轉折 4：input 成為「重構者」")
math(s,["dx/dt = F( x(t), u(t) )"], y=2.4, size=30)
idea(s,["輸入不只是訊號：initial condition x₀ 與 input u(t)","共同指定一個運算。"],"purple",y=4.2,size=22)
note(s,"reservoir「丟石頭進池塘」→ 為特定運算量身打造的動力系統 · Driscoll/Golub/Sussillo 2018")

# 20 why hands-on
s=new(); title(s,"我為什麼自己跑一遍", kicker="PART 4", kc="brown")
center_big(s,["光看論文不夠。用 numpy + scipy 寫最小可運行版本，","圖跟數字都是真實 RLS 計算出來的——零捏造。"],base=22)
cards(s,[("① sinusoid","FORCE-Output","green"),("② Lorenz","FORCE-Internal v3","brown"),
         ("③ fixed-point","解剖 trained 網路","purple")], y=4.6, h=1.7)

# 21 exp1 image
s=new(); title(s,"實驗一：sinusoid（FORCE-Output）")
image(s,"force_sinusoid_demo.png","N=1000, g=1.5, 15 train cycles — 紅線整段覆蓋灰線，test error ±0.008 內。")

# 22 exp1 read
s=new(); title(s,"判讀：誤差從第一步就小")
bullets(s,["**Cycle 1   MSE = 0.000019    |W| = 0.093",
           "Cycle 5   MSE = 0.000048    |W| = 0.196",
           "Cycle 15  MSE = 0.000005    |W| = 0.251",
           "**Cycle 16 (test)  MSE = 0.000006"], size=22)
idea(s,["Cycle 1 就是 1.9×10⁻⁵ — 式 (7)(8) 的實證。"],"green",y=5.4,size=22)

# 23 exp2 image
s=new(); title(s,"實驗二：Lorenz（FORCE-Internal v3）")
image(s,"force_internal_v3_demo.png","左：時序前段貼合、後段 phase divergence。右：delay-embedding 重建 Lorenz 蝴蝶。")

# 24 exp2 read
s=new(); title(s,"判讀：混沌被保留")
cards(s,[("|J|_norm: 0.168 → 0.169","全程僅變 0.7% → chaotic regime 完整保留","brown"),
         ("運算靠 W 承載","|W|: 0.09 → 0.54",None)], y=2.6, h=2.2)
note(s,"test MSE 在 cycle 32+ 衝到 0.2–0.5 是混沌 phase divergence 必然 — 判讀看 attractor，不看 point-by-point。")

# 25 exp3 image
s=new(); title(s,"實驗三：fixed-point 分析（Sussillo & Barak 2013 風格）")
image(s,"stage4_fixed_points_demo.png","4-panel：PC1-PC2 / PC1-PC3 / Jacobian 特徵值雲 / q-landscape。123 個 slow points。")

# 26 exp3 read manifold
s=new(); title(s,"判讀：低維 manifold")
idea(s,["top-3 PC 解釋 ≈ 81% 變異（55.2% + 19.2% + 6.6%）"],"purple",y=2.6,size=24)
center_big(s,["高維 RNN 的運算落在低維流形上；","slow points 聚在兩個 lobe 中心 → 對應 Lorenz 的 C⁺/C⁻。"],base=22)

# 27 exp3 saddles
s=new(); title(s,"判讀：unstable saddles 驅動 lobe-switch")
bullets(s,["特徵值雲多數在 Re≈−1（因 F=−x+… 主導）",
           "但有可觀部分跨過虛軸 Re=0 進入 unstable 區",
           "強 saddle 最多 26 個 unstable directions, max Re≈0.40",
           "**沒有 q<10⁻⁷ 的真不動點 — 這是 feature"], size=21)
note(s,"RNN 用連續的 slow manifold（而非離散吸引子）在「計算」混沌。")

# 28 CP3
s=new(); vote(s,"RNN 找不到真正的 fixed point（q<10⁻⁷），代表什麼？",
  ["A. 訓練失敗","B. 網路太小","C. 它用連續的 slow manifold 編碼混沌","D. 演算法 bug"],
  "C","答案 C：離散答案是錯問題；slow manifold 在 N 維空間是連續的。")

# 29 closed loop
s=new(); title(s,"三實驗閉環")
pivot_card(s,0.7,2.6,3.7,1.6,"FORCE 訓練","轉折 1","","green","")
pivot_card(s,4.8,2.6,3.7,1.6,"fixed-point 解剖","轉折 2","","purple","")
pivot_card(s,8.9,2.6,3.7,1.6,"低維 manifold","CTPD","微縮實證","gold","")
center_big(s,["轉折 2 的方法，作用在轉折 1 的網路上，產出 CTPD 框架的縮影。"],base=22)

# 30 CTPD
s=new(); title(s,"終點：CTPD 框架", kicker="PART 5", kc="gold")
idea(s,["Computation through Neural Population Dynamics","Vyas, Golub, Sussillo & Shenoy 2020, Annu. Rev. Neurosci."],"gold",size=24)
cards(s,[("Motor control","",None),("Timing","",None),("Decision","",None),("Working memory","",None)], y=4.6, h=1.6)

# 31 equation
s=new(); title(s,"一條式子貫穿全場")
math(s,["dx/dt = f( x(t), u(t) )"], y=2.6, size=34)
center_big(s,["x 是 N 維神經群體狀態。神經群體就是一個動力系統，","運算透過它在 state space 的時間演化來實現。"],base=22)

# 32 critique
s=new(); title(s,"批判：守住兩條界線")
cards(s,[("RNN ≠ 大腦","只是「運算可以如何被解」的一個範例，非 ground truth。是 hypothesis-generating tool。","brown"),
         ("低維投影的陷阱","2D/3D 投影幾乎可 fit 任何假設 → 必須在高維（10–20 維）量化。","brown")], y=2.6, h=3.0)

# 33 emergence
s=new(); title(s,"哲學底色：Emergence")
idea(s,["Conway's Game of Life：簡單規則 → 複雜湧現"],None,y=2.6,size=26)
note(s,"Sussillo 2026 Q&A 自述為終生母題。", y=4.0)
center_big(s,["運算從大量簡單單元的集體動力中湧現，","而非定位於單一元件 — 即 population dynamics 的信念。"],base=20)

# 34 takehome
s=new(); title(s,"帶走三句話")
bullets(s,["FORCE 用 P 矩陣在第一步把誤差打到近乎零 → 保留並馴服混沌。",
           "真正生出 CTPD 的，是把 RNN 當動力系統來解剖（fixed points）。",
           "高維神經運算落在低維 manifold，由 fixed/slow points 的幾何決定。"], size=22)

# 35 refs
s=new(); title(s,"延伸閱讀")
bullets(s,["Sussillo & Abbott 2009 · Neuron — FORCE learning",
           "Sussillo & Barak 2013 · Neural Computation — Opening the Black Box",
           "Sussillo 2014 · Curr. Opin. Neurobiol. — Circuits as dynamical systems",
           "Driscoll, Golub & Sussillo 2018 · Neuron — Computation through Cortical Dynamics",
           "Vyas, Golub, Sussillo & Shenoy 2020 · Annu. Rev. Neurosci. — CTPD"], size=19)
note(s,"三張圖皆由 numpy/scipy 實際計算產生（基於 Sussillo & Abbott 2009 最小重現），零捏造。")

out = HERE/"deck.pptx"
prs.save(str(out))
print("OK ->", out, "slides:", len(prs.slides._sldIdLst))
